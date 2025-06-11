from fastapi import FastAPI, UploadFile, File, Form
from fastapi.responses import StreamingResponse, JSONResponse
from pptx import Presentation
import io
import os
import requests

app = FastAPI(title="PowerPoint Customizer")

OPEN_ROUTER_API_KEY = os.getenv("OPEN_ROUTER_API_KEY")
OPEN_ROUTER_MODEL = os.getenv("OPEN_ROUTER_MODEL", "mistralai/mistral-7b-instruct")


def call_openrouter(original_text: str, context: str) -> str:
    """Send text to OpenRouter and return the generated result."""
    if not OPEN_ROUTER_API_KEY:
        return original_text  # no API key provided, return original

    payload = {
        "model": OPEN_ROUTER_MODEL,
        "messages": [
            {"role": "system", "content": f"Use this context to customize the presentation: {context}"},
            {"role": "user", "content": original_text},
        ],
    }
    headers = {
        "Authorization": f"Bearer {OPEN_ROUTER_API_KEY}",
        "Content-Type": "application/json",
    }
    try:
        resp = requests.post(
            "https://openrouter.ai/api/v1/chat/completions",
            json=payload,
            headers=headers,
            timeout=30,
        )
        resp.raise_for_status()
        data = resp.json()
        return data["choices"][0]["message"]["content"]
    except Exception:
        return original_text


def generate_customized_pptx(prs: Presentation, context: str) -> Presentation:
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                updated_text = call_openrouter(shape.text, context)
                shape.text = updated_text
    return prs


@app.post("/preview")
async def preview_presentation(pptx_file: UploadFile = File(...), context: str = Form(...)):
    contents = await pptx_file.read()
    prs = Presentation(io.BytesIO(contents))
    preview = []
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                new_text = call_openrouter(shape.text, context)
                preview.append({
                    "slide": slide_idx,
                    "shape": shape_idx,
                    "original": shape.text,
                    "generated": new_text,
                })
    return JSONResponse(preview)


@app.post("/apply")
async def apply_changes(pptx_file: UploadFile = File(...), context: str = Form(...)):
    contents = await pptx_file.read()
    prs = Presentation(io.BytesIO(contents))
    prs = generate_customized_pptx(prs, context)
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=customized.pptx"},
    )

